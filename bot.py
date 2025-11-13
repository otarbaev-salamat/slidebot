import asyncio
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, FSInputFile
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.client.session.aiohttp import AiohttpSession
from pathlib import Path
from pptx import Presentation
from translateapi import translate_batch  # sizning alohida translate.py
import os

BOT_TOKEN = "8508767861:AAGvTqzevWzCIicsIGJkzHeBQqFxNLK6Bk4"

FILES_DIR = Path("files")
FILES_DIR.mkdir(exist_ok=True)

session = AiohttpSession(timeout=900)
bot = Bot(token=BOT_TOKEN, session=session)
storage = MemoryStorage()
dp = Dispatcher(storage=storage)


class TranslateStates(StatesGroup):
    waiting_for_file = State()
    waiting_for_source_language = State()
    waiting_for_target_language = State()


LANGUAGE_CODES = {
    "Qaraqalpaqsha": "kaa_Latn",
    "Қарақалпақша": "kaa_Cyrl",
    "O‘zbekcha": "uzn_Latn",
    "Ўзбекча": "uzn_Cyrl",
    "Qazaqsha": "kaz_Cyrl",
    "Русски": "rus_Cyrl",
    "Inglizsha": "eng_Latn",
    "Turkshe": "tur_Latn",
    "Koreysshe": "kor_Hang",
    "Xitaysha": "zho_Hans",
}


# === HANDLERS === #
@dp.message(Command("start"))
async def start(message: types.Message, state: FSMContext):
    await state.clear()
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✨ Awdarmalaw", callback_data="translate_slide")],
        [InlineKeyboardButton(text="ℹ️ Jardem", callback_data="help")]
    ])
    await message.answer("🎉 Xosh keldin'iz! Fayldi awdarmalaw ushın tanlań:", reply_markup=keyboard)


@dp.callback_query(F.data == "translate_slide")
async def ask_file(callback: types.CallbackQuery, state: FSMContext):
    await callback.answer()
    await callback.message.answer("📂 Iltimas, .pptx faylin jiberin:")
    await state.set_state(TranslateStates.waiting_for_file)


@dp.message(TranslateStates.waiting_for_file, F.document)
async def handle_file(message: types.Message, state: FSMContext):
    document = message.document
    if not document.file_name.lower().endswith(".pptx"):
        await message.answer("❌ Tekgana .pptx fayllarni jiberin.")
        return

    await state.update_data(file_id=document.file_id, file_name=document.file_name)
    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [InlineKeyboardButton(text=name, callback_data=f"src:{code}")]
            for name, code in LANGUAGE_CODES.items()
        ]
    )
    await message.answer("📘 Tilin tanlań:", reply_markup=kb)
    await state.set_state(TranslateStates.waiting_for_source_language)


@dp.callback_query(TranslateStates.waiting_for_source_language, F.data.startswith("src:"))
async def choose_source_lang(callback: types.CallbackQuery, state: FSMContext):
    await state.update_data(source_lang=callback.data.split(":")[1])
    kb = InlineKeyboardMarkup(
        inline_keyboard=[
            [InlineKeyboardButton(text=name, callback_data=f"tgt:{code}")]
            for name, code in LANGUAGE_CODES.items()
        ]
    )
    await callback.message.edit_text("🎯 Qaysi tilge awdarmalaysiz:", reply_markup=kb)
    await state.set_state(TranslateStates.waiting_for_target_language)


@dp.callback_query(TranslateStates.waiting_for_target_language, F.data.startswith("tgt:"))
async def process_translation(callback: types.CallbackQuery, state: FSMContext):
    await callback.answer()
    data = await state.get_data()
    file_id = data["file_id"]
    file_name = data["file_name"]  # original nom
    source_lang = data["source_lang"]
    target_lang = callback.data.split(":")[1]

    status_msg = await callback.message.edit_text("📥 Fayl jüklenbekte...")
    input_path = FILES_DIR / file_name
    temp_output_path = FILES_DIR / f"translated_{file_name}"  # vaqtinchalik tarjima fayli

    # Faylni yuklash
    file = await bot.get_file(file_id)
    await bot.download_file(file.file_path, destination=input_path)

    await status_msg.edit_text("⚙️ Awdarmalaw baslandi...")
    await translate_pptx(input_path, temp_output_path, source_lang, target_lang, status_msg)

    # Foydalanuvchiga original nomi bilan jo‘natish
    await status_msg.edit_text("✅ Awdarmalaw tamamlandi!")
    await callback.message.answer_document(FSInputFile(temp_output_path), caption="Faylin'iz tayin ✅", filename=file_name)

    # Fayllarni tozalash
    try:
        os.remove(input_path)
        os.remove(temp_output_path)
    except Exception as e:
        print(f"Faylni o‘chirishda xato: {e}")

    await state.clear()


# === TRANSLATE PPTX === #
async def translate_pptx(input_path, output_path, source_lang, target_lang, status_msg):
    prs = Presentation(input_path)

    def collect_all_texts(shape, collected):
        if shape.shape_type == 6 and hasattr(shape, "shapes"):  # GROUP
            for s in shape.shapes:
                collect_all_texts(s, collected)
        elif shape.shape_type == 19:  # TABLE
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame") and cell.text_frame:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                if r.text.strip():
                                    collected.append((r, r.text))
        elif hasattr(shape, "text_frame") and shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        collected.append((r, r.text))

    # Barcha matnlarni yig‘ish
    all_texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            collect_all_texts(shape, all_texts)

    total = len(all_texts)
    if total == 0:
        await status_msg.edit_text("⚠️ Slaydda tekst tabilmadi.")
        return

    # Tarjima (batch async)
    batch_size = 100
    translated_results = []
    for i in range(0, total, batch_size):
        chunk = [t[1] for t in all_texts[i:i+batch_size]]
        # translate.py dagi translate_batch ishlatamiz
        results = await asyncio.to_thread(translate_batch, chunk, source_lang, target_lang)
        translated_results.extend(results)
        await status_msg.edit_text(f"⏳ Awdarmalaw: {min(i+batch_size, total)}/{total}")

    # Natijani yozish
    for (run, _), translated_text in zip(all_texts, translated_results):
        run.text = translated_text

    prs.save(output_path)


# === RUN === #
async def main():
    print("✅ Bot ishga tushdi...")
    await dp.start_polling(bot, skip_updates=True)


if __name__ == "__main__":
    asyncio.run(main())
