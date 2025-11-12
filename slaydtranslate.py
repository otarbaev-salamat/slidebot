from pptx import Presentation
from translateapi import translate
from concurrent.futures import ThreadPoolExecutor, as_completed
import shutil
import time
from tqdm import tqdm
from pptx.dml.color import RGBColor
import traceback

def slaydt(input_pptx, output_pptx, source_lang, target_lang):

    log_file = output_pptx.replace('.pptx', '_log.txt')

    shutil.copyfile(input_pptx, output_pptx)
    prs = Presentation(output_pptx)

    def extract_shapes(shape, collected):
        if shape.shape_type == 6:
            for s in shape.shapes:
                extract_shapes(s, collected)
        elif shape.shape_type == 19:
            for row in shape.table.rows:
                for cell in row.cells:
                    collected.append(cell)
        elif hasattr(shape, "text_frame") and shape.has_text_frame:
            collected.append(shape)

    all_text_elements = []
    for slide in prs.slides:
        for shape in slide.shapes:
            extract_shapes(shape, all_text_elements)

    print(f"🔹 Tabilgan elementlar sani: {len(all_text_elements)}")

    paragraphs = []
    for shape in all_text_elements:
        tf = shape.text_frame
        for p in tf.paragraphs:
            text = " ".join(run.text for run in p.runs if run.text.strip()).strip()
            if text:
                paragraphs.append((p, text))

    print(f"🔸 Tabilgan paragraphlar sani: {len(paragraphs)}")

    def translate_safe(text, source_language, target_language):
        try:
            result = translate(text, source_lang=source_language, target_lang=target_language)
            if isinstance(result, str):
                return result
            elif isinstance(result, dict) and "sentences" in result:
                translated = " ".join(
                    s["translated"] for s in result["sentences"] if "translated" in s
                )
                return translated
            else:
                return None
        except Exception as e:
            log_entry(text, "", "TRANSLATE FAILED ❌", traceback.format_exc())
            return None

    def log_entry(original, translated, status, error_message=""):
        with open(log_file, "a", encoding="utf-8") as f:
            f.write(f"\n---\nSTATUS: {status}\nERROR: {error_message}\nORIGINAL: {original}\nTRANSLATED: {translated}\n")

    translated_texts = []
    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = {executor.submit(translate_safe, text, source_lang, target_lang): (p, text) for p, text in paragraphs}
        for future in tqdm(as_completed(futures), total=len(futures), desc="Translating"):
            p, text = futures[future]
            translated = future.result()
            if translated:
                translated_texts.append((p, translated))
            else:
                log_entry(text, "", "FAILED ⚠️", "Tarjima xatosi")

    for p, translated in translated_texts:
        try:
            if p.runs:
                ref_run = p.runs[0]
                font = ref_run.font
                bold = font.bold
                italic = font.italic
                size = font.size
                name = font.name

                try:
                    if font.color and font.color.rgb:
                        color = font.color.rgb
                    else:
                        color = None
                except Exception as e:
                    log_entry(p.text, translated, "COLOR ERROR 🎨", str(e))
                    color = None
            else:
                ref_run = None

            p.clear()
            new_run = p.add_run()
            if translated:
                new_run.text = translated
            else:
                new_run.text = text

            if ref_run:
                f = new_run.font
                f.bold = bold
                f.italic = italic
                f.size = size
                f.name = name

                if color:
                    try:
                        f.color.rgb = color
                    except Exception as e:
                        log_entry(p.text, translated, "WRITE COLOR ERROR 🎨", str(e))
                        f.color.rgb = RGBColor(0, 0, 0)
        except Exception as e:
            log_entry(p.text, translated, f"WRITE ERROR ❌", traceback.format_exc())

    prs.save(output_pptx)
    return f"\n💾 Saved: {output_pptx}"
# slaydt('test.pptx','tested.pptx','uzn_Latn','kaa_Latn')